"""Build a client-ready PowerPoint from Section3_Q9_2022_Client_Presentation.md."""

from __future__ import annotations

import re
from dataclasses import dataclass, field
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

REPO_ROOT = Path(__file__).resolve().parents[1]
MD_PATH = REPO_ROOT / "docs" / "Section3_Q9_2022_Client_Presentation.md"
OUTPUT_PATH = REPO_ROOT / "docs" / "Section3_Q9_2022_Client_Presentation.pptx"
ASSETS_ROOT = REPO_ROOT / "docs"

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
MARGIN_L = Inches(0.5)
MARGIN_R = Inches(0.5)
CONTENT_TOP = Inches(1.05)
CONTENT_BOTTOM = Inches(6.85)
CONTENT_WIDTH = SLIDE_WIDTH - MARGIN_L - MARGIN_R
FOOTER_TOP = Inches(7.02)
GAP = Inches(0.1)

COLOR_NAVY = RGBColor(31, 58, 95)
COLOR_ACCENT = RGBColor(46, 107, 158)
COLOR_TEXT = RGBColor(51, 51, 51)
COLOR_MUTED = RGBColor(102, 102, 102)
COLOR_WHITE = RGBColor(255, 255, 255)


@dataclass
class SlideContent:
    number: int
    title: str
    blocks: list[dict] = field(default_factory=list)


def strip_md_inline(text: str) -> str:
    text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
    return text.replace("`", "").strip()


def is_table_separator(line: str) -> bool:
    if not line.strip().startswith("|"):
        return False
    cells = [cell.strip() for cell in line.strip().strip("|").split("|")]
    if not cells:
        return False
    return all(re.fullmatch(r":?-{3,}:?", cell) for cell in cells)


def merge_adjacent_tables(blocks: list[dict]) -> list[dict]:
    merged: list[dict] = []
    for block in blocks:
        if block["type"] == "table" and merged and merged[-1]["type"] == "table":
            merged[-1]["rows"].extend(block["rows"])
        else:
            merged.append(block)
    return merged


def parse_markdown(path: Path) -> list[SlideContent]:
    raw = path.read_text(encoding="utf-8")
    parts = re.split(r"\n## Slide (\d+): ", raw)
    preamble = parts[0].strip()
    slides: list[SlideContent] = []
    if preamble:
        title = preamble.lstrip("#").split("\n")[0].strip()
        slides.append(SlideContent(number=0, title=title))

    for idx in range(1, len(parts), 2):
        number = int(parts[idx])
        body = parts[idx + 1]
        first_line, _, rest = body.partition("\n")
        slide = SlideContent(number=number, title=first_line.strip())
        table_rows: list[list[str]] = []
        in_table = False

        for line in rest.splitlines():
            stripped = line.strip()
            if not stripped:
                if in_table and table_rows:
                    slide.blocks.append({"type": "table", "rows": table_rows})
                    table_rows = []
                    in_table = False
                continue

            if stripped.startswith("### "):
                if in_table and table_rows:
                    slide.blocks.append({"type": "table", "rows": table_rows})
                    table_rows = []
                    in_table = False
                slide.blocks.append({"type": "heading", "text": strip_md_inline(stripped[4:])})
                continue

            if stripped.startswith("!["):
                if in_table and table_rows:
                    slide.blocks.append({"type": "table", "rows": table_rows})
                    table_rows = []
                    in_table = False
                match = re.match(r"!\[(.*?)\]\((.*?)\)", stripped)
                if match:
                    slide.blocks.append({"type": "image", "alt": match.group(1), "path": match.group(2)})
                continue

            if stripped.startswith("|"):
                if is_table_separator(stripped):
                    continue
                cells = [strip_md_inline(c) for c in stripped.strip("|").split("|")]
                table_rows.append(cells)
                in_table = True
                continue

            if in_table and table_rows:
                slide.blocks.append({"type": "table", "rows": table_rows})
                table_rows = []
                in_table = False

            if stripped.startswith("- "):
                slide.blocks.append({"type": "bullet", "text": strip_md_inline(stripped[2:])})
                continue

            if stripped.startswith(">"):
                slide.blocks.append({"type": "quote", "text": strip_md_inline(stripped.lstrip("> ").strip())})
                continue

            slide.blocks.append({"type": "paragraph", "text": strip_md_inline(stripped)})

        if in_table and table_rows:
            slide.blocks.append({"type": "table", "rows": table_rows})

        slide.blocks = merge_adjacent_tables(slide.blocks)
        slides.append(slide)

    return slides


def parse_pairwise_bullet(text: str) -> dict[str, str] | None:
    """Parse a slide bullet into the fields shown in the pairwise comparison table."""

    match = re.match(
        r"Black Hills vs (?P<benchmark>.+?):\s*"
        r"(?P<direction>.+?);\s*"
        r"(?P<significance>significant|not significant);\s*"
        r"Black Hills median = (?P<bh_median>.+?);\s*"
        r".+? median = (?P<benchmark_median>.+?);\s*"
        r"median gap = (?P<median_gap>.+?);\s*"
        r"p\s*(?P<p_value>.+?)\s*\.?\s*$",
        text,
        flags=re.IGNORECASE,
    )
    if not match:
        return None
    return {
        "benchmark": match.group("benchmark").strip(),
        "direction": match.group("direction").strip(),
        "significance": match.group("significance").strip().lower(),
        "bh_median": match.group("bh_median").strip(),
        "benchmark_median": match.group("benchmark_median").strip(),
        "median_gap": match.group("median_gap").strip(),
        "p_value": match.group("p_value").strip(),
    }


def estimate_text_height(text: str, font_size: int, width_inches: float = 12.2) -> float:
    chars_per_line = max(24, int(width_inches * 72 / (font_size * 0.55)))
    lines = max(1, (len(text) + chars_per_line - 1) // chars_per_line)
    return max(0.28, lines * (font_size / 72) * 1.35)


class SlideRenderer:
    def __init__(self, prs: Presentation) -> None:
        self.prs = prs
        self.slide = None
        self.y = CONTENT_TOP
        self.slide_num = 0
        self.slide_title = ""
        self.subtitle = ""

    def new_slide(self, slide_num: int, title: str, subtitle: str | None = None) -> None:
        layout = self.prs.slide_layouts[6]
        self.slide = self.prs.slides.add_slide(layout)
        self.slide_num = slide_num
        self.slide_title = title
        self.subtitle = subtitle or ""
        self.y = CONTENT_TOP

        bar = self.slide.shapes.add_shape(1, 0, 0, SLIDE_WIDTH, Inches(0.92))
        bar.fill.solid()
        bar.fill.fore_color.rgb = COLOR_NAVY
        bar.line.fill.background()

        title_h = Inches(0.5) if subtitle else Inches(0.55)
        box = self.slide.shapes.add_textbox(MARGIN_L, Inches(0.15), CONTENT_WIDTH, title_h)
        tf = box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(20 if subtitle else 22 if slide_num else 28)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE
        if subtitle:
            p2 = tf.add_paragraph()
            p2.text = subtitle
            p2.font.size = Pt(10)
            p2.font.color.rgb = RGBColor(210, 220, 235)

        footer = self.slide.shapes.add_textbox(MARGIN_L, FOOTER_TOP, CONTENT_WIDTH, Inches(0.3))
        fp = footer.text_frame.paragraphs[0]
        footer_text = f"Q9 2022 · {slide_num}" if slide_num else ""
        if self.subtitle:
            footer_text = f"{footer_text} · {self.subtitle}"
        fp.text = footer_text
        fp.font.size = Pt(9)
        fp.font.color.rgb = COLOR_MUTED
        fp.alignment = PP_ALIGN.RIGHT

    def _remaining(self) -> float:
        return CONTENT_BOTTOM - self.y

    def _need_continuation(self, needed_inches: float = 0.5) -> bool:
        return self._remaining() < Inches(needed_inches)

    def continue_slide(self) -> None:
        subtitle = f"{self.subtitle} (continued)" if self.subtitle else "(continued)"
        self.new_slide(self.slide_num, self.slide_title, subtitle)

    def add_paragraph(self, text: str, font_size: int = 11, *, bold: bool = False, italic: bool = False) -> None:
        h = Inches(estimate_text_height(text, font_size))
        h = min(h, self._remaining())
        if h <= 0:
            return
        box = self.slide.shapes.add_textbox(MARGIN_L, self.y, CONTENT_WIDTH, h)
        tf = box.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.vertical_anchor = MSO_ANCHOR.TOP
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.italic = italic
        p.font.color.rgb = COLOR_TEXT
        self.y += h + GAP

    def add_bullet(self, text: str, font_size: int = 10) -> None:
        h = Inches(estimate_text_height(f"• {text}", font_size))
        h = min(h, self._remaining())
        if h <= 0:
            return
        box = self.slide.shapes.add_textbox(MARGIN_L + Inches(0.12), self.y, CONTENT_WIDTH - Inches(0.12), h)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"• {text}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = COLOR_TEXT
        self.y += h + GAP

    def add_heading(self, text: str) -> None:
        self.add_paragraph(text, 12, bold=True)

    def add_table(self, rows: list[list[str]], font_size: int = 9) -> None:
        if not rows:
            return
        row_h_in = 0.26 if font_size <= 8 else 0.3
        needed = Inches(len(rows) * row_h_in)
        if self._remaining() < needed:
            header, *body = rows[0], rows[1:]
            rows_per_slide = max(1, int(self._remaining() / Inches(row_h_in)) - 1)
            for start in range(0, len(body), rows_per_slide):
                if self._need_continuation(row_h_in * 2):
                    self.continue_slide()
                chunk = [header, *body[start : start + rows_per_slide]]
                self.add_table(chunk, font_size=font_size)
            return

        nrows = len(rows)
        ncols = max(len(r) for r in rows)
        height = Inches(nrows * row_h_in)
        if height <= 0:
            return
        shape = self.slide.shapes.add_table(nrows, ncols, MARGIN_L, self.y, CONTENT_WIDTH, height)
        table = shape.table
        for r_idx, row in enumerate(rows):
            for c_idx in range(ncols):
                cell = table.cell(r_idx, c_idx)
                cell.text = row[c_idx] if c_idx < len(row) else ""
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                cell.margin_left = cell.margin_right = Inches(0.04)
                cell.margin_top = cell.margin_bottom = Inches(0.02)
                for para in cell.text_frame.paragraphs:
                    para.font.size = Pt(font_size)
                    para.font.color.rgb = COLOR_WHITE if r_idx == 0 else COLOR_TEXT
                    if r_idx == 0:
                        para.font.bold = True
                if r_idx == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = COLOR_ACCENT
        self.y = shape.top + shape.height + GAP

    def add_table_at(
        self,
        left: int,
        top: int,
        width: int,
        rows: list[list[str]],
        *,
        font_size: int = 8,
        row_h_in: float = 0.19,
    ) -> int:
        """Place a table at fixed position; returns bottom edge (no slide splits)."""
        if not rows:
            return top
        nrows = len(rows)
        ncols = max(len(r) for r in rows)
        height = Inches(nrows * row_h_in)
        shape = self.slide.shapes.add_table(nrows, ncols, left, top, width, height)
        table = shape.table
        if ncols >= 2:
            table.columns[0].width = Inches(1.55)
            table.columns[1].width = width - Inches(1.55)
        for r_idx, row in enumerate(rows):
            for c_idx in range(ncols):
                cell = table.cell(r_idx, c_idx)
                cell.text = row[c_idx] if c_idx < len(row) else ""
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                cell.margin_left = cell.margin_right = Inches(0.03)
                cell.margin_top = cell.margin_bottom = Inches(0.01)
                for para in cell.text_frame.paragraphs:
                    para.font.size = Pt(font_size)
                    para.font.color.rgb = COLOR_WHITE if r_idx == 0 else COLOR_TEXT
                    if r_idx == 0:
                        para.font.bold = True
                if r_idx == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = COLOR_ACCENT
        return int(shape.top + shape.height)

    def add_text_at(
        self,
        left: int,
        top: int,
        width: int,
        text: str,
        font_size: int,
        *,
        bold: bool = False,
        bullet: bool = False,
    ) -> int:
        """Place a text box at fixed position; returns bottom edge."""
        display = f"• {text}" if bullet else text
        width_in = width / Inches(1)
        h_in = estimate_text_height(display, font_size, width_inches=width_in)
        height = Inches(h_in)
        box = self.slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.vertical_anchor = MSO_ANCHOR.TOP
        p = tf.paragraphs[0]
        p.text = display
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = COLOR_TEXT
        gap = int(GAP / 4) if bullet else int(GAP / 2)
        return int(box.top + box.height + gap)

    def add_pairwise_table(self, bullets: list[str]) -> None:
        rows = [parse_pairwise_bullet(t) for t in bullets]
        rows = [r for r in rows if r]
        if not rows:
            for t in bullets:
                self.add_bullet(t, 9)
            return
        has_gap = any(r.get("median_gap") for r in rows)
        header = ["Benchmark", "Significance", "BH median", "Benchmark median"]
        if has_gap:
            header.append("Median gap (95% CI)")
        header.append("p-value")
        table: list[list[str]] = [header]
        for r in rows:
            sig = "Significant" if r["significance"] == "significant" else "Not significant"
            row_cells = [r["benchmark"], sig, r["bh_median"], r["benchmark_median"]]
            if has_gap:
                row_cells.append(r.get("median_gap") or "—")
            row_cells.append(f"p {r['p_value']}")
            table.append(row_cells)
        self.add_table(table, font_size=8 if has_gap else 9)

    def add_image(
        self,
        path: Path,
        *,
        max_height: float | None = None,
        reserve_below_in: float = 0.0,
        fill_available: bool = False,
    ) -> None:
        if not path.exists():
            return
        if self._need_continuation(1.0):
            self.continue_slide()
        remaining_in = self._remaining() / Inches(1.0)
        if fill_available:
            max_h_in = remaining_in - reserve_below_in - 0.12
        elif max_height is not None:
            max_h_in = min(max_height, remaining_in - 0.12)
        else:
            max_h_in = remaining_in - 0.12
        if max_h_in <= 0.55:
            return
        max_h = Inches(max_h_in)
        pic = self.slide.shapes.add_picture(str(path), MARGIN_L, self.y, width=CONTENT_WIDTH)
        if pic.height > max_h:
            scale = max_h / pic.height
            pic.height = int(max_h)
            pic.width = int(pic.width * scale)
            pic.left = int(MARGIN_L + (CONTENT_WIDTH - pic.width) / 2)
        self.y = pic.top + pic.height + GAP

    def estimate_blocks_height_in(self, blocks: list[dict], *, font_size: int = 10) -> float:
        total = 0.0
        for block in blocks:
            if block["type"] == "paragraph":
                total += estimate_text_height(block["text"], font_size) + 0.1
            elif block["type"] == "bullet":
                total += estimate_text_height(f"• {block['text']}", font_size) + 0.08
            elif block["type"] == "heading":
                total += 0.32
            elif block["type"] == "table":
                fs = 8 if len(block["rows"]) > 8 else 9
                row_h = 0.26 if fs <= 8 else 0.3
                total += len(block["rows"]) * row_h + 0.12
            elif block["type"] == "pairwise_table":
                total += 5 * 0.26 + 0.12
            elif block["type"] == "quote":
                total += estimate_text_height(block["text"], 11) + 0.1
        return total


def group_blocks_for_chart(blocks: list[dict]) -> list[dict]:
    """Reorder chart-slide blocks: caption, image, pairwise table."""
    image = next((b for b in blocks if b["type"] == "image"), None)
    caption = next(
        (b for b in blocks if b["type"] == "paragraph" and b["text"].lower().startswith("pairwise")),
        None,
    )
    bullets = [b["text"] for b in blocks if b["type"] == "bullet"]
    other = [
        b
        for b in blocks
        if b not in ([image] if image else [])
        and b != caption
        and b["type"] not in ("bullet",)
    ]
    ordered: list[dict] = []
    if caption:
        ordered.append(caption)
    if image:
        ordered.append(image)
    if bullets:
        ordered.append({"type": "pairwise_table", "bullets": bullets})
    ordered.extend(other)
    return ordered


def render_chart_slide(renderer: SlideRenderer, blocks: list[dict]) -> None:
    """Caption, large chart, comparison table at bottom."""
    blocks = group_blocks_for_chart(blocks)
    caption = next((b for b in blocks if b["type"] == "paragraph"), None)
    image = next((b for b in blocks if b["type"] == "image"), None)
    table_block = next((b for b in blocks if b["type"] == "pairwise_table"), None)
    tail = [b for b in blocks if b["type"] not in ("paragraph", "image", "pairwise_table")]

    if caption:
        renderer.add_paragraph(caption["text"], 9)

    reserve = renderer.estimate_blocks_height_in(tail, font_size=10)
    if table_block:
        reserve += 5 * 0.28 + 0.12

    if image:
        renderer.add_image(
            ASSETS_ROOT / image["path"],
            reserve_below_in=reserve,
            fill_available=True,
        )
    if table_block:
        renderer.add_pairwise_table(table_block["bullets"])
    if tail:
        render_blocks(renderer, tail, font_size=10)


def render_overview_slide(renderer: SlideRenderer, blocks: list[dict]) -> None:
    """Slide 2: large overview figure with bullets underneath."""
    image = next((b for b in blocks if b["type"] == "image"), None)
    rest = [b for b in blocks if b["type"] != "image"]
    reserve = renderer.estimate_blocks_height_in(rest, font_size=10) + 0.1
    if image:
        renderer.add_image(ASSETS_ROOT / image["path"], reserve_below_in=reserve, fill_available=True)
    render_blocks(renderer, rest, font_size=10)


def render_blocks(
    renderer: SlideRenderer,
    blocks: list[dict],
    *,
    chart_mode: bool = False,
    font_size: int = 10,
    bullet_font_size: int | None = None,
) -> None:
    bullet_font_size = bullet_font_size or font_size

    for block in blocks:
        if renderer._need_continuation(0.25):
            renderer.continue_slide()
        if block["type"] == "paragraph":
            renderer.add_paragraph(block["text"], font_size)
        elif block["type"] == "heading":
            renderer.add_heading(block["text"])
        elif block["type"] == "bullet":
            renderer.add_bullet(block["text"], bullet_font_size)
        elif block["type"] == "quote":
            renderer.add_paragraph(block["text"], font_size + 1, italic=True)
        elif block["type"] == "table":
            fs = 8 if len(block["rows"]) > 8 else 9
            renderer.add_table(block["rows"], font_size=fs)
        elif block["type"] == "image":
            renderer.add_image(ASSETS_ROOT / block["path"], max_height=5.0, fill_available=True)
        elif block["type"] == "pairwise_table":
            renderer.add_pairwise_table(block["bullets"])


def _split_blocks_by_heading(slide: SlideContent) -> list[tuple[str, list[dict]]]:
    """Split slide blocks into (subtitle, blocks) sections at ### headings."""
    sections: list[tuple[str, list[dict]]] = []
    subtitle = "Overview"
    current: list[dict] = []
    for block in slide.blocks:
        if block["type"] == "heading":
            if current:
                sections.append((subtitle, current))
            subtitle = block["text"]
            current = []
            continue
        current.append(block)
    if current:
        sections.append((subtitle, current))
    return sections


def build_slide_1(prs: Presentation, slide: SlideContent) -> None:
    """One slide: dataset on top, variable definitions underneath (full width)."""
    sections = {name: blocks for name, blocks in _split_blocks_by_heading(slide)}
    dataset_blocks = sections.get("Dataset", [])
    variable_blocks = sections.get("Variable definitions", [])

    content_w = int(CONTENT_WIDTH)
    x = int(MARGIN_L)
    y = int(CONTENT_TOP)

    r = SlideRenderer(prs)
    r.new_slide(1, slide.title)

    y = r.add_text_at(x, y, content_w, "Dataset", 11, bold=True)
    for block in dataset_blocks:
        if block["type"] == "bullet":
            y = r.add_text_at(x, y, content_w, block["text"], 8, bullet=True)
        elif block["type"] == "paragraph":
            y = r.add_text_at(x, y, content_w, block["text"], 8)

    y = r.add_text_at(x, y + int(GAP), content_w, "Variable definitions", 11, bold=True)
    for block in variable_blocks:
        if block["type"] == "paragraph":
            y = r.add_text_at(x, y, content_w, block["text"], 8)
        elif block["type"] == "table":
            y = r.add_table_at(x, y, content_w, block["rows"], font_size=7, row_h_in=0.17)
            y += int(GAP / 2)


def build_slide_14(prs: Presentation, slide: SlideContent) -> None:
    """Two slides: region counts + headline; change table + notes."""
    split_idx = None
    for i, b in enumerate(slide.blocks):
        if b["type"] == "heading" and "significance or direction" in b["text"].lower():
            split_idx = i
            break
    first = slide.blocks[:split_idx] if split_idx else slide.blocks
    second = slide.blocks[split_idx:] if split_idx else []

    r1 = SlideRenderer(prs)
    r1.new_slide(slide.number, slide.title, "Excluded organizations & headline")
    render_blocks(r1, first, font_size=11, bullet_font_size=11)

    if second:
        r2 = SlideRenderer(prs)
        r2.new_slide(slide.number, slide.title, "What changed when special orgs are included")
        render_blocks(r2, second, font_size=11, bullet_font_size=11)


def build_slide_15(prs: Presentation, slide: SlideContent) -> None:
    sections: list[tuple[str, list[dict]]] = [
        ("Summary & methods", []),
        ("Significant results", []),
        ("Board conclusions", []),
    ]
    current = 0
    for block in slide.blocks:
        if block["type"] == "heading":
            lower = block["text"].lower()
            if "statistically significant pairwise" in lower:
                current = 1
                sections[1][1].append(block)
                continue
            if "no significant pairwise" in lower or "how to state" in lower or "one-sentence" in lower:
                current = 2
                sections[2][1].append(block)
                continue
        sections[current][1].append(block)

    subtitles = [s[0] for s in sections]
    for subtitle, blocks in zip(subtitles, [s[1] for s in sections]):
        if not blocks:
            continue
        r = SlideRenderer(prs)
        r.new_slide(slide.number, slide.title, subtitle)
        render_blocks(r, blocks, font_size=11, bullet_font_size=11)


def build_presentation() -> Path:
    slides = parse_markdown(MD_PATH)
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    for slide in slides:
        if slide.number == 0:
            r = SlideRenderer(prs)
            r.new_slide(0, "Q9 2022 Client Presentation")
            box = r.slide.shapes.add_textbox(Inches(1.2), Inches(2.5), Inches(10.9), Inches(1.8))
            tf = box.text_frame
            p = tf.paragraphs[0]
            p.text = "Q9 2022 Client Presentation"
            p.font.size = Pt(34)
            p.font.bold = True
            p.font.color.rgb = COLOR_NAVY
            p.alignment = PP_ALIGN.CENTER
            p2 = tf.add_paragraph()
            p2.text = "Revenue sources — Black Hills vs benchmark regions"
            p2.font.size = Pt(17)
            p2.font.color.rgb = COLOR_ACCENT
            p2.alignment = PP_ALIGN.CENTER
            continue

        if slide.number == 1:
            build_slide_1(prs, slide)
            continue
        if slide.number == 13:
            build_slide_14(prs, slide)
            continue
        if slide.number == 14:
            build_slide_15(prs, slide)
            continue

        r = SlideRenderer(prs)
        r.new_slide(slide.number, slide.title)
        if slide.number == 2:
            render_overview_slide(r, slide.blocks)
        elif 3 <= slide.number <= 12:
            render_chart_slide(r, slide.blocks)
        else:
            render_blocks(r, slide.blocks, font_size=11, bullet_font_size=11)

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    try:
        prs.save(str(OUTPUT_PATH))
        return OUTPUT_PATH
    except PermissionError:
        fallback = OUTPUT_PATH.with_stem(OUTPUT_PATH.stem + "_updated")
        prs.save(str(fallback))
        print(
            f"Warning: could not overwrite {OUTPUT_PATH.name} (file may be open). "
            f"Saved to {fallback.name} instead."
        )
        return fallback


if __name__ == "__main__":
    print(f"Saved: {build_presentation()}")
